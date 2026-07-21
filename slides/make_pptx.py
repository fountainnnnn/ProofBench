"""Generate the ProofBench hackathon pitch deck as a .pptx.

Mirrors slides/index.html: same DESIGN.md tokens (OKLCH converted to sRGB),
same copy, 12 slides, 16:9. Tooling only, not part of the app.

Run: .venv/Scripts/python slides/make_pptx.py
Out: slides/proofbench_pitch.pptx
"""

import math
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- tokens

def oklch(l, c, h):
    """OKLCH -> sRGB hex (gamut-clamped)."""
    hr = math.radians(h)
    a, b = c * math.cos(hr), c * math.sin(hr)
    l_ = l + 0.3963377774 * a + 0.2158037573 * b
    m_ = l - 0.1055613458 * a - 0.0638541728 * b
    s_ = l - 0.0894841775 * a - 1.2914855480 * b
    l3, m3, s3 = l_ ** 3, m_ ** 3, s_ ** 3
    r = +4.0767416621 * l3 - 3.3077115913 * m3 + 0.2309699292 * s3
    g = -1.2684380046 * l3 + 2.6097574011 * m3 - 0.3413193965 * s3
    bl = -0.0041960863 * l3 - 0.7034186147 * m3 + 1.7076147010 * s3

    def gam(v):
        v = 12.92 * v if v <= 0.0031308 else 1.055 * v ** (1 / 2.4) - 0.055
        return round(min(max(v, 0), 1) * 255)

    return RGBColor(gam(r), gam(g), gam(bl))

BG = oklch(0.992, 0.002, 264)
SURFACE = oklch(0.985, 0.003, 264)
SURFACE_2 = oklch(0.972, 0.004, 264)
BORDER = oklch(0.905, 0.006, 264)
BORDER_STRONG = oklch(0.82, 0.008, 264)
TEXT = oklch(0.24, 0.02, 264)
TEXT_2 = oklch(0.46, 0.015, 264)
TEXT_3 = oklch(0.62, 0.012, 264)
ACCENT = oklch(0.52, 0.19, 268)
ACCENT_SOFT = oklch(0.95, 0.025, 268)
CODE_BG = oklch(0.27, 0.02, 264)
CODE_TEXT = oklch(0.85, 0.02, 160)
CODE_DIM = oklch(0.65, 0.02, 200)
T_INFO = oklch(0.78, 0.12, 230)
T_WARN = oklch(0.80, 0.14, 80)
T_OK = oklch(0.80, 0.14, 155)
T_ERR = oklch(0.75, 0.15, 25)
WHITE = oklch(0.99, 0.005, 268)

SANS = "Inter"
MONO = "Consolas"

PAGE_W, PAGE_H = 13.333, 7.5
MARGIN = 0.92
CONTENT_W = PAGE_W - 2 * MARGIN

# Services used in the hackathon build. Names identify integrations only; they
# do not imply sponsorship, endorsement, partnership, or affiliation.
HACKATHON_INTEGRATIONS = ["Daytona", "Oxylabs", "Doubleword"]

NO_ENDORSEMENT_NOTE = (
    "Third-party names and trademarks identify integrations only and do not imply "
    "sponsorship, endorsement, partnership, or affiliation."
)

# ---------------------------------------------------------------- helpers

def _style(run, text, st):
    run.text = text
    f = run.font
    f.name = st.get("font", SANS)
    f.size = Pt(st.get("size", 18))
    f.bold = st.get("bold", False)
    f.color.rgb = st.get("color", TEXT)
    if st.get("tracking"):
        f._rPr.set("spc", str(st["tracking"]))


def add_text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts with runs=[(text, style)], plus optional
    space_before/space_after (pt), line_spacing, align."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_before"):
            para.space_before = Pt(p["space_before"])
        if p.get("space_after"):
            para.space_after = Pt(p["space_after"])
        if p.get("line_spacing"):
            para.line_spacing = p["line_spacing"]
        for text, st in p["runs"]:
            _style(para.add_run(), text, st)
    return tb


def para(runs, **kw):
    d = {"runs": runs}
    d.update(kw)
    return d


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, round_=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ is not None else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if round_ is not None:
        try:
            sp.adjustments[0] = round_
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    sp.text_frame.margin_left = sp.text_frame.margin_right = 0
    sp.text_frame.margin_top = sp.text_frame.margin_bottom = 0
    return sp


def hline(slide, x, y, w, color=BORDER):
    add_rect(slide, x, y, w, 0.012, fill=color)


def chrome(slide, idx, total):
    add_rect(slide, 0, 0, PAGE_W * (idx + 1) / total, 0.045, fill=ACCENT)
    add_text(slide, MARGIN, PAGE_H - 0.42, 2.0, 0.3,
             [para([("ProofBench", {"size": 10.5, "bold": True, "color": TEXT_2})])])
    add_text(slide, PAGE_W - MARGIN - 1.5, PAGE_H - 0.42, 1.5, 0.3,
             [para([(f"{idx + 1:02d} / {total:02d}", {"size": 10, "font": MONO, "color": TEXT_3})],
                   align=PP_ALIGN.RIGHT)])


def numbered_rows(slide, rows, y0=2.55, row_h=1.06, num_size=14, title_size=18.5, desc_size=16):
    y = y0
    for num, title, desc in rows:
        hline(slide, MARGIN, y, CONTENT_W)
        add_text(slide, MARGIN, y + 0.18, 0.62, 0.4,
                 [para([(num, {"size": num_size, "font": MONO, "bold": True, "color": ACCENT})])])
        add_text(slide, MARGIN + 0.85, y + 0.14, CONTENT_W - 0.85, row_h - 0.2, [
            para([(title, {"size": title_size, "bold": True, "color": TEXT})], space_after=3),
            para([(desc, {"size": desc_size, "color": TEXT_2})], line_spacing=1.15),
        ])
        y += row_h
    hline(slide, MARGIN, y, CONTENT_W)
    return y


def integrations_line(slide, y, prefix=None):
    names = (["Built with"] if prefix else []) + HACKATHON_INTEGRATIONS
    runs = []
    for i, n in enumerate(names):
        runs.append((n + ("        " if i < len(names) - 1 else ""),
                     {"size": 13, "color": TEXT_3, "bold": bool(prefix) and i == 0}))
    add_text(slide, MARGIN, y, CONTENT_W, 0.35, [para(runs)])


def h2(slide, text, y=0.85, size=40, color=TEXT, w=CONTENT_W):
    add_text(slide, MARGIN, y, w, 1.4,
             [para([(text, {"size": size, "bold": True, "color": color, "tracking": -60})],
                   line_spacing=1.04)])


# ---------------------------------------------------------------- deck

prs = Presentation()
prs.slide_width = Emu(Inches(PAGE_W))
prs.slide_height = Emu(Inches(PAGE_H))

# Deck attribution is the product name, not a legal owner. "ProofBench" here is
# the product this deck is about; it names no company, entity, or rights holder.
prs.core_properties.title = "ProofBench"
prs.core_properties.author = "ProofBench"
prs.core_properties.last_modified_by = "ProofBench"
# Fixed timestamps keep regeneration byte-stable; the default is python-pptx's
# 2013 template date. `comments` is dc:description in core.xml.
prs.core_properties.created = datetime(2026, 7, 20, 0, 0, 0)
prs.core_properties.modified = datetime(2026, 7, 20, 0, 0, 0)
prs.core_properties.subject = "ProofBench hackathon pitch deck"
prs.core_properties.comments = (
    "Product attribution only; third-party names identify integrations and do "
    "not imply endorsement."
)
BLANK = prs.slide_layouts[6]
TOTAL = 12
idx = 0


def new_slide():
    global idx
    s = prs.slides.add_slide(BLANK)
    # Native slide background instead of a full-canvas shape: a PAGE_W x PAGE_H
    # rectangle reads as an overflowing shape to slide QA tooling.
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    chrome(s, idx, TOTAL)
    idx += 1
    return s


# 01 · Title
s = new_slide()
add_text(s, MARGIN, 1.95, 11.0, 1.9,
         [para([("ProofBench", {"size": 88, "bold": True, "color": TEXT, "tracking": -320})])])
add_text(s, MARGIN, 3.75, 9.9, 1.9, [
    para([("Agentic benchmarking of any tool, built from its own docs. ",
           {"size": 21, "bold": True, "color": TEXT}),
          ("Point it at a category and a labelled dataset: candidates are discovered, "
           "integrated, and run in parallel sandboxes, then scored deterministically against "
           "your ground truth.", {"size": 21, "color": TEXT_2})], line_spacing=1.3),
])
integrations_line(s, 5.85)

# 02 · Problem
s = new_slide()
h2(s, "Picking a tool is still a guess.")
numbered_rows(s, [
    ("01", "Marketing over measurement.",
     "Every vendor page claims best-in-class performance. None of them show numbers on your data."),
    ("02", "Public benchmarks do not transfer.",
     "Leaderboards run on generic datasets. Your workload is not in them."),
    ("03", "A real eval costs a week.",
     "Every team hand-builds the same integrations, reruns the same harness, then argues about how to score it."),
], y0=2.7, row_h=1.35)

# 03 · Relevance
s = new_slide()
add_text(s, MARGIN, 1.45, 11.4, 1.1,
         [para([("Every stack is a chain of tool choices: OCR, parsers, vector stores, inference APIs.",
                 {"size": 25, "color": TEXT_2})], line_spacing=1.25)])
add_text(s, MARGIN, 2.95, 10.6, 1.7,
         [para([("A few points of quality become manual review hours. Latency and cost per "
                 "1,000 calls decide what ships.", {"size": 34, "bold": True, "color": TEXT})],
               line_spacing=1.2)])
add_text(s, MARGIN, 5.1, 11.2, 1.6,
         [para([("The wrong pick compounds. Teams still choose on vibes.",
                 {"size": 44, "bold": True, "color": ACCENT, "tracking": -40})], line_spacing=1.15)])

# 04 · What it does
s = new_slide()
h2(s, "Point it at a category and your labelled dataset.")
add_text(s, MARGIN, 2.15, 11.2, 1.3, [
    para([("ProofBench returns a ", {"size": 20, "color": TEXT_2}),
          ("ranked, citation-backed verdict", {"size": 20, "bold": True, "color": TEXT}),
          (" on which tool actually performs best on your data. No LLM ever judges "
           "correctness; the LLM only does logistics.", {"size": 20, "color": TEXT_2})],
         line_spacing=1.3),
])
steps = [
    ("1", "Discover", "Search for candidate tools in the category.", "orchestrator · Oxylabs"),
    ("2", "Read the docs", "Scrape documentation and pricing pages.", "Oxylabs"),
    ("3", "Integrate", "Generate a working adapter per tool, from its own docs.", "Doubleword"),
    ("4", "Run", "Build and execute every candidate in parallel.", "Daytona"),
    ("5", "Score", "Deterministic metrics against your ground truth.", "evaluate.py"),
]
col_w = CONTENT_W / 5
y0 = 4.15
for k, (n, name, desc, who) in enumerate(steps):
    x = MARGIN + k * col_w
    if k:
        add_rect(s, x, y0 + 0.05, 0.012, 2.0, fill=BORDER)
    add_text(s, x + (0.28 if k else 0), y0, col_w - 0.42, 2.3, [
        para([(n, {"size": 12, "font": MONO, "bold": True, "color": ACCENT})], space_after=8),
        para([(name, {"size": 18.5, "bold": True, "color": TEXT})], space_after=6),
        para([(desc, {"size": 13, "color": TEXT_2})], line_spacing=1.25, space_after=10),
        para([(who, {"size": 11, "font": MONO, "color": TEXT_3})]),
    ])

# 05 · Hackathon integrations
s = new_slide()
h2(s, "Three integrations, one pipeline.")
tbl_y = 2.3
add_rect(s, MARGIN, tbl_y, CONTENT_W, 0.5, fill=SURFACE_2)
add_text(s, MARGIN + 0.2, tbl_y + 0.11, 2.4, 0.3,
         [para([("Integration", {"size": 12, "bold": True, "color": TEXT_2})])])
add_text(s, MARGIN + 2.8, tbl_y + 0.11, 8.0, 0.3,
         [para([("Role in ProofBench", {"size": 12, "bold": True, "color": TEXT_2})])])
integration_rows = [
    ("Daytona", "Sandbox fleet. ",
     "A pre-warmed pool of isolated sandboxes. One per candidate, built and run in parallel.", 0.85),
    ("Oxylabs", "Docs intel. ",
     "Google search for candidate discovery, scraping of documentation and pricing pages.", 0.85),
    ("Doubleword", "Batch inference and codegen. ",
     "Generates each tool's adapter from its scraped documentation, and batches every LLM request "
     "in the pipeline: codegen calls and hosted extraction candidates all run through Doubleword's "
     "batch serving. It also competes in the benchmark as a candidate itself.", 1.35),
]
y = tbl_y + 0.5
for name, lead, rest, rh in integration_rows:
    hline(s, MARGIN, y, CONTENT_W)
    add_text(s, MARGIN + 0.2, y + 0.17, 2.5, 0.5,
             [para([(name, {"size": 17.5, "bold": True, "color": TEXT})])])
    add_text(s, MARGIN + 2.8, y + 0.17, CONTENT_W - 3.0, rh - 0.25,
             [para([(lead, {"size": 16, "bold": True, "color": TEXT}),
                    (rest, {"size": 16, "color": TEXT_2})], line_spacing=1.18)])
    y += rh
hline(s, MARGIN, y, CONTENT_W)
add_text(s, MARGIN, y + 0.22, CONTENT_W, 0.4,
         [para([(NO_ENDORSEMENT_NOTE, {"size": 11, "color": TEXT_3})], line_spacing=1.2)])

# 06 · Feature 1: docs intel
s = new_slide()
h2(s, "It reads the docs so you do not have to.", y=1.55, size=36, w=6.1)
add_text(s, MARGIN, 3.35, 6.1, 3.2, [
    para([("Give the orchestrator agent a plain-language brief: 'benchmark invoice-extraction "
           "tools', or any other category. It searches Google and scrapes documentation and "
           "pricing pages with Oxylabs, then proposes an ", {"size": 18, "color": TEXT_2}),
          ("editable spec", {"size": 18, "bold": True, "color": TEXT}),
          (": candidates, dataset, metrics. You confirm, it runs.", {"size": 18, "color": TEXT_2})],
         line_spacing=1.35),
])
card_x, card_y, card_w, card_h = 7.45, 1.55, 4.95, 4.8
add_rect(s, card_x, card_y, card_w, card_h, fill=SURFACE, line=BORDER, round_=0.045)
add_text(s, card_x + 0.32, card_y + 0.28, card_w - 0.6, 0.3,
         [para([("Proposed spec", {"size": 12, "bold": True, "color": TEXT_2})])])
chips = [("tesseract", True), ("easyocr", True), ("paddleocr", True),
         ("doubleword", True), ("+ discovered tools", False)]
cx, cy = card_x + 0.32, card_y + 0.72
for label, tinted in chips:
    w = 0.34 + len(label) * 0.088
    if cx + w > card_x + card_w - 0.3:
        cx = card_x + 0.32
        cy += 0.46
    pill = add_rect(s, cx, cy, w, 0.36, fill=ACCENT_SOFT if tinted else SURFACE_2,
                    line=BORDER, round_=0.5)
    tf = pill.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _style(p.add_run(), label, {"size": 12, "font": MONO,
                                "color": TEXT if tinted else TEXT_2})
    cx += w + 0.14
hline(s, card_x + 0.32, card_y + 2.35, card_w - 0.64)
add_text(s, card_x + 0.32, card_y + 2.55, card_w - 0.64, 1.4, [
    para([("Dataset: ", {"size": 13, "color": TEXT_2}), ("data/demo", {"size": 13, "bold": True, "color": TEXT}),
          (", 15 labelled invoices", {"size": 13, "color": TEXT_2})], space_after=5),
    para([("Metrics: ", {"size": 13, "color": TEXT_2}),
          ("exact accuracy, field F1, CER, latency, cost", {"size": 13, "bold": True, "color": TEXT})],
         space_after=5),
    para([("Demo category: ", {"size": 13, "color": TEXT_2}),
          ("invoice extraction", {"size": 13, "bold": True, "color": TEXT})]),
])
btn = add_rect(s, card_x + 0.32, card_y + 4.0, 1.85, 0.46, fill=ACCENT, round_=0.16)
btf = btn.text_frame
btf.word_wrap = False
bp = btf.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
_style(bp.add_run(), "Run benchmark", {"size": 13.5, "bold": True, "color": WHITE})

# 07 · Feature 2: Doubleword codegen + batch inference
s = new_slide()
h2(s, "Integrations written from each tool's own documentation.", y=1.4, size=36, w=6.1)
add_text(s, MARGIN, 3.15, 6.1, 3.6, [
    para([("Doubleword reads the scraped docs and generates the full adapter: install commands, "
           "API calls, and output mapping to ", {"size": 17.5, "color": TEXT_2}),
          ("one shared contract", {"size": 17.5, "bold": True, "color": TEXT}),
          (". The sandbox builds and validates it; a failed build goes back for regeneration.",
           {"size": 17.5, "color": TEXT_2})], line_spacing=1.3),
    para([("The same Doubleword deployment is also the ", {"size": 17.5, "color": TEXT_2}),
          ("batch inference layer", {"size": 17.5, "bold": True, "color": TEXT}),
          (": every LLM request in a run, from adapter codegen to the hosted extraction "
           "candidates, goes out as a batched request through Doubleword.",
           {"size": 17.5, "color": TEXT_2})], line_spacing=1.3, space_before=14),
])
code_x, code_y, code_w, code_h = 7.45, 1.5, 4.95, 3.9
add_rect(s, code_x, code_y, code_w, code_h, fill=SURFACE, line=BORDER, round_=0.03)
COMMENT = {"size": 12.5, "font": MONO, "color": TEXT_3}
PLAIN = {"size": 12.5, "font": MONO, "color": TEXT}
KEY = {"size": 12.5, "font": MONO, "color": ACCENT, "bold": True}
code_lines = [
    [("# easyocr adapter (from scraped docs)", COMMENT)],
    [("def", KEY), (" extract(image_path):", PLAIN)],
    [("    reader = easyocr.Reader([\"en\"])", PLAIN)],
    [("    lines = reader.readtext(image_path)", PLAIN)],
    [("    return", KEY), (" {", PLAIN)],
    [("        \"invoice_no\": parse_no(lines),", PLAIN)],
    [("        \"date\":       parse_date(lines),", PLAIN)],
    [("        \"total\":      parse_total(lines),", PLAIN)],
    [("    }", PLAIN)],
]
add_text(s, code_x + 0.32, code_y + 0.3, code_w - 0.6, code_h - 0.5,
         [para([(t, st) for t, st in line], line_spacing=1.45) for line in code_lines])

# 08 · Feature 3: sandbox fleet
s = new_slide()
h2(s, "Every candidate races in its own Daytona sandbox.", y=1.55, size=36, w=6.1)
add_text(s, MARGIN, 4.0, 6.1, 2.6, [
    para([("A pre-warmed pool builds, validates, and runs candidates ", {"size": 18.5, "color": TEXT_2}),
          ("in parallel", {"size": 18.5, "bold": True, "color": TEXT}),
          (". A failed build triggers self-repair: the agent reads the error, patches the "
           "adapter, retries. Every line streams to the console.", {"size": 18.5, "color": TEXT_2})],
         line_spacing=1.35),
])
term_x, term_y, term_w, term_h = 7.45, 1.6, 4.95, 3.35
add_rect(s, term_x, term_y, term_w, term_h, fill=CODE_BG, round_=0.045)
term_lines = [
    ("[sb-1]", " pip install easyocr torch ", "building", T_INFO),
    ("[sb-1]", " build ok in 41s ", "validating", T_WARN),
    ("[sb-2]", " tesseract: exit 1 ", "failed", T_ERR),
    ("[sb-2]", " repair 1: adapter patched ", "repair", T_WARN),
    ("[sb-2]", " rebuild ok ", "running", T_INFO),
    ("[sb-1]", " 15/15 docs extracted ", "done", T_OK),
]
add_text(s, term_x + 0.32, term_y + 0.28, term_w - 0.6, term_h - 0.5, [
    para([(prefix, {"size": 12.5, "font": MONO, "color": CODE_DIM}),
          (body, {"size": 12.5, "font": MONO, "color": CODE_TEXT}),
          (status, {"size": 12.5, "font": MONO, "color": sc})], line_spacing=1.5)
    for prefix, body, status, sc in term_lines
])

# 09 · Feature 4: deterministic evaluator
s = new_slide()
h2(s, "Scored by code, never judged by a model.", y=1.35, size=36, w=6.1)
add_text(s, MARGIN, 3.3, 6.1, 1.2, [
    para([("evaluate.py", {"size": 17.5, "bold": True, "font": MONO, "color": TEXT}),
          (" compares every result against your ground truth. Deterministic, offline, "
           "reproducible byte-for-byte. The metric set adapts to the category.",
           {"size": 17.5, "color": TEXT_2})], line_spacing=1.3),
])
metric_rows = [
    ("exact accuracy", "whole-document match"),
    ("field F1 · CER", "per-field and character-level error"),
    ("latency · failure rate · cost/1k", "the operational picture"),
]
my = 4.7
for name, desc in metric_rows:
    hline(s, MARGIN, my, 6.1)
    add_text(s, MARGIN, my + 0.12, 6.1, 0.4,
             [para([(name + "   ", {"size": 14.5, "font": MONO, "bold": True, "color": TEXT}),
                    (desc, {"size": 15.5, "color": TEXT_2})])])
    my += 0.52
hline(s, MARGIN, my, 6.1)

res_x, res_y, res_w = 7.15, 1.35, 5.35
add_rect(s, res_x, res_y, res_w, 0.36 + 4 * 0.44 + 0.06, fill=SURFACE, line=BORDER, round_=0.03)
cols = [("#", 0.35), ("Candidate", 1.22), ("Exact", 0.62), ("F1", 0.55),
        ("CER", 0.55), ("Latency", 0.68), ("Fail", 0.50), ("$/1k", 0.60)]
add_rect(s, res_x + 0.02, res_y + 0.02, res_w - 0.04, 0.36, fill=SURFACE_2)
cx = res_x + 0.16
for label, cw in cols:
    add_text(s, cx, res_y + 0.09, cw - 0.04, 0.26,
             [para([(label, {"size": 10, "bold": True, "color": TEXT_2})],
                   align=PP_ALIGN.LEFT if label == "Candidate" else PP_ALIGN.RIGHT)], wrap=False)
    cx += cw
res_rows = [
    ("1", "openai_vision", "96.7%", "97.5%", "0.012", "1.18s", "0.0%", "$7.50"),
    ("2", "doubleword", "94.2%", "95.6%", "0.021", "2.41s", "0.0%", "$2.10"),
    ("3", "easyocr", "82.5%", "86.1%", "0.084", "0.74s", "3.3%", "$0.00"),
    ("4", "tesseract", "77.5%", "81.4%", "0.112", "0.31s", "6.7%", "$0.00"),
]
ry = res_y + 0.38
for ri, row in enumerate(res_rows):
    if ri == 0:
        add_rect(s, res_x + 0.02, ry, res_w - 0.04, 0.44, fill=ACCENT_SOFT)
    else:
        hline(s, res_x + 0.14, ry, res_w - 0.28)
    cx = res_x + 0.16
    for ci, (val, (label, cw)) in enumerate(zip(row, cols)):
        style = {"size": 11, "font": MONO, "color": TEXT if ri == 0 else TEXT_2, "bold": ri == 0}
        if label == "Candidate":
            style = {"size": 11.5, "bold": True, "color": TEXT}
        add_text(s, cx, ry + 0.11, cw - 0.04, 0.3,
                 [para([(val, style)], align=PP_ALIGN.LEFT if label == "Candidate" else PP_ALIGN.RIGHT)],
                 wrap=False)
        cx += cw
    ry += 0.44
add_text(s, res_x, ry + 0.12, res_w + 0.3, 0.3,
         [para([("Sample output, invoice-extraction demo category. One of many possible tool categories.",
                 {"size": 11, "color": TEXT_3})])])

# 10 · Feature 5: console
s = new_slide()
h2(s, "The agent's work is on stage.")
numbered_rows(s, [
    ("01", "Chat console.",
     "Launch a benchmark in one sentence. The chat is a console, not a chatbot."),
    ("02", "Live trace feed.",
     "Every tool call, every docs scrape, every repair attempt, streamed over SSE as it happens."),
    ("03", "Per-sandbox terminals.",
     "Watch each Daytona sandbox build and run, in the open."),
    ("04", "Ranked report.",
     "Sortable results table plus a citation-backed markdown report, downloadable as PDF. "
     "It never invents a number."),
], y0=2.5, row_h=1.12)

# 11 · Demo
s = new_slide()
h2(s, "The three-minute demo.")
numbered_rows(s, [
    ("1", "Say what to benchmark.", "\u201cBenchmark invoice-extraction tools on my invoices.\u201d"),
    ("2", "Review the spec.", "The agent searches and scrapes docs, then proposes candidates."),
    ("3", "Attach data.", "The synthetic set of 15 labelled invoices, or upload your own."),
    ("4", "Run.", "Four Daytona sandboxes build and race in parallel, repairs live on screen."),
    ("5", "Read the verdict.", "Ranked deterministic metrics and a citation-backed report."),
], y0=2.45, row_h=0.88, num_size=13, title_size=17, desc_size=15)

# 12 · Close
s = new_slide()
add_text(s, MARGIN, 2.1, 12.0, 1.6,
         [para([("Proven, not promised.", {"size": 66, "bold": True, "color": TEXT, "tracking": -160})])])
add_text(s, MARGIN, 3.85, 10.9, 1.7, [
    para([("ProofBench tells you which tool wins on your data, with numbers you can defend. ",
           {"size": 19, "bold": True, "color": TEXT}),
          ("OCR is the demo, not the limit: any tool category with public docs qualifies, "
           "and the same loop becomes a regression gate in CI.", {"size": 19, "color": TEXT_2})],
         line_spacing=1.35),
])
integrations_line(s, 5.75, prefix=True)
add_text(s, MARGIN, 6.25, CONTENT_W, 0.4,
         [para([(NO_ENDORSEMENT_NOTE, {"size": 11, "color": TEXT_3})], line_spacing=1.2)])

# ---------------------------------------------------------------- save

out = Path(__file__).with_name("proofbench_pitch.pptx")
prs.save(out)
print(f"wrote {out} ({TOTAL} slides)")
